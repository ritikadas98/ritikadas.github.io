# -*- coding: utf-8 -*-
"""Pack a folder of slide images into a 16:9 .pptx.

    python build.py <image-dir> <out.pptx> ["Deck title"]

Each slide is one full-bleed picture. The text is not editable in PowerPoint,
which is the trade: what you get back is exactly what the web deck looks like,
including the gradients, the charts and the diagrams, rather than an
approximation rebuilt in PowerPoint's own type engine.

Images are re-encoded to JPEG on the way in. A 3840px PNG of a dark gradient is
about 2 MB; the same frame at quality 92 is a fifth of that, and nobody can see
the difference at projector size.
"""
import io
import os
import sys
import glob

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

WIDE_W, WIDE_H = Inches(13.333), Inches(7.5)   # 16:9
JPEG_QUALITY = 92


def build(img_dir, out_path, title=None):
    shots = sorted(glob.glob(os.path.join(img_dir, 'slide-*.png')))
    if not shots:
        sys.exit('no slide-*.png found in %s' % img_dir)

    prs = Presentation()
    prs.slide_width, prs.slide_height = WIDE_W, WIDE_H
    blank = prs.slide_layouts[6]

    for shot in shots:
        slide = prs.slides.add_slide(blank)
        with Image.open(shot) as im:
            buf = io.BytesIO()
            im.convert('RGB').save(buf, format='JPEG', quality=JPEG_QUALITY, optimize=True)
            buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, width=WIDE_W, height=WIDE_H)

    if title:
        prs.core_properties.title = title

    os.makedirs(os.path.dirname(out_path) or '.', exist_ok=True)
    prs.save(out_path)
    size = os.path.getsize(out_path) / (1024 * 1024)
    print('%s — %d slides, %.1f MB' % (out_path, len(shots), size))


if __name__ == '__main__':
    if len(sys.argv) < 3:
        sys.exit('usage: python build.py <image-dir> <out.pptx> ["Deck title"]')
    build(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else None)
