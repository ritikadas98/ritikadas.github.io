"""Deck engine — one layout system, a palette per project.

Every case study gets its own colour and the same six layouts, so the decks read
as a set without reading as copies. The palettes are not invented here: they are
lifted from tools/og-card.py, which has carried a brand and a four-stop field per
project since the cards were first built. A reader who saw the card recognises
the deck.

16:9 at 13.333 x 7.5in, which is what Google Slides expects on import.
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W, H = 13.333, 7.5
DISPLAY, BODY, MONO = "Bricolage Grotesque 14pt", "Inter", "IBM Plex Mono"

def hx(s): return C.from_string(s.lstrip("#"))

class Deck:
    def __init__(self, brand, field, accent, paper="#EDE8E1", ink="#211E24", ink2="#4A4550"):
        self.brand, self.field, self.accent = brand, field, accent
        self.paper, self.ink, self.ink2 = paper, ink, ink2
        self.p = Presentation()
        self.p.slide_width, self.p.slide_height = In(W), In(H)

    # ---- primitives ----
    def _blank(self, bg=None):
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(W), In(H))
        r.fill.solid(); r.fill.fore_color.rgb = hx(bg or self.paper); r.line.fill.background()
        r.shadow.inherit = False
        return s

    def _text(self, s, x, y, w, h, runs, size=18, font=BODY, colour=None, bold=False,
              align=PP_ALIGN.LEFT, spacing=1.15, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        items = runs if isinstance(runs, list) else [runs]
        for i, item in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align; para.line_spacing = spacing
            txt, opts = (item, {}) if isinstance(item, str) else item
            r = para.add_run(); r.text = txt
            f = r.font
            f.name = opts.get("font", font); f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.color.rgb = hx(opts.get("colour", colour or self.ink))
            if opts.get("space_after"): para.space_after = Pt(opts["space_after"])
        return tb

    def _card(self, s, x, y, w, h, fill=None, line=None, radius=0.06):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
        c.adjustments[0] = radius
        c.fill.solid(); c.fill.fore_color.rgb = hx(fill or "#FFFFFF")
        if line: c.line.color.rgb = hx(line); c.line.width = Pt(1)
        else: c.line.fill.background()
        c.shadow.inherit = False
        return c

    def _pill(self, s, label, x=0.75, y=0.5):
        w = 0.28 + 0.105 * len(label)
        p = self._card(s, x, y, w, 0.42, fill=self.brand, radius=0.5)
        self._text(s, x, y + 0.085, w, 0.3, label.upper(), size=11, font=MONO,
                   colour="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        return p

    def _source(self, s, text):
        self._text(s, 0.75, H - 0.72, W - 1.5, 0.3, text, size=10.5, font=BODY,
                   colour="#8A8188")

    # ---- the six layouts ----
    def cover(self, kicker, title, sub, byline):
        """Full-bleed field gradient. The only slide that uses the brand at full strength."""
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(W), In(H))
        bg.line.fill.background(); bg.shadow.inherit = False
        bg.fill.gradient()
        stops = bg.fill.gradient_stops
        stops[0].color.rgb = hx(self.field[0]); stops[0].position = 0.0
        stops[1].color.rgb = hx(self.field[3]); stops[1].position = 1.0
        bg.fill.gradient_angle = 45.0
        self._text(s, 0.9, 0.6, 8, 0.4, kicker.upper(), size=11, font=MONO, colour="#FFFFFF", bold=True)
        self._text(s, 0.9, 2.15, 9.6, 2.4, title, size=54, font=DISPLAY, colour="#FFFFFF", bold=True, spacing=1.0)
        self._text(s, 0.9, 4.55, 8.2, 1.2, sub, size=19, colour="#FFFFFF", spacing=1.35)
        self._text(s, 0.9, H - 1.35, 8, 0.8, byline, size=13, colour="#FFFFFF", spacing=1.35)
        return s

    def split(self, pill, title, body, rows=None, source=None):
        """Heading and prose left, labelled rows right. The workhorse."""
        s = self._blank(); self._pill(s, pill)
        self._text(s, 0.75, 1.5, 5.4, 2.6, title, size=36, font=DISPLAY, bold=True, spacing=1.05)
        self._text(s, 0.75, 4.2, 5.2, 2.2, body, size=15, colour=self.ink2, spacing=1.5)
        y = 1.5
        for k, v in (rows or []):
            self._card(s, 6.75, y, 5.8, 1.28)
            self._text(s, 7.05, y + 0.2, 5.2, 0.28, k.upper(), size=10, font=MONO, colour=self.brand, bold=True)
            self._text(s, 7.05, y + 0.54, 5.2, 0.62, v, size=12.5, colour=self.ink2, spacing=1.35)
            y += 1.42
        if source: self._source(s, source)
        return s

    def cards(self, pill, title, sub, items, source=None):
        """Two or four white cards. For anything that is genuinely a set."""
        s = self._blank(); self._pill(s, pill)
        self._text(s, 0.75, 1.45, 11.8, 0.9, title, size=38, font=DISPLAY, bold=True, spacing=1.05)
        if sub: self._text(s, 0.75, 2.45, 10.5, 0.5, sub, size=15, colour=self.ink2, spacing=1.4)
        n = len(items); cols = 2 if n <= 4 else 3
        cw = (11.8 - 0.4 * (cols - 1)) / cols; ch = 2.02 if n <= 4 else 1.72
        for i, (k, v) in enumerate(items):
            x = 0.75 + (i % cols) * (cw + 0.4); y = 3.2 + (i // cols) * (ch + 0.35)
            self._card(s, x, y, cw, ch)
            self._text(s, x + 0.35, y + 0.35, cw - 0.7, 0.4, k, size=17, font=DISPLAY, bold=True, colour=self.brand)
            self._text(s, x + 0.35, y + 0.88, cw - 0.7, ch - 1.05, v, size=12, colour=self.ink2, spacing=1.45)
        if source: self._source(s, source)
        return s

    def compare(self, pill, title, sub, left, right, source=None):
        """Grey card for the state without, brand card for the state with."""
        s = self._blank(); self._pill(s, pill)
        self._text(s, 0.75, 1.45, 11.8, 1.5, title, size=38, font=DISPLAY, bold=True, spacing=1.05)
        if sub: self._text(s, 0.75, 2.9, 10.5, 0.5, sub, size=15, colour=self.ink2, spacing=1.4)
        for i, (head, lines, fill, fg) in enumerate([
                (left[0], left[1], "#8C8A88", "#FFFFFF"), (right[0], right[1], self.brand, "#FFFFFF")]):
            x = 0.75 + i * 6.1
            self._card(s, x, 3.7, 5.7, 2.6, fill=fill)
            self._text(s, x + 0.4, 3.95, 5, 0.3, head.upper(), size=10, font=MONO, colour=fg, bold=True)
            self._text(s, x + 0.4, 4.4, 5, 1.7,
                       [(l, {"size": 20 if j == 0 else 13, "font": DISPLAY if j == 0 else BODY,
                             "bold": j == 0, "colour": fg, "space_after": 8}) for j, l in enumerate(lines)],
                       spacing=1.35)
        if source: self._source(s, source)
        return s

    def spine(self, pill, title, sub, steps, highlight=None, source=None):
        """Numbered rows, head and note in fixed columns.

        Two earlier attempts failed the same way. Stacking the note under the head
        needs ~0.8in a row and six rows will not fit under a title, so the rows
        collided. Putting both in one text flow let the head push the note onto a
        second line, which then overflowed a short card.

        Fixed columns fix both: the head owns 2.7in, the note owns the rest, and
        neither can push the other anywhere.
        """
        s = self._blank(); self._pill(s, pill)
        self._text(s, 0.75, 1.35, 11.8, 0.8, title, size=36, font=DISPLAY, bold=True, spacing=1.05)
        if sub: self._text(s, 0.75, 2.25, 11, 0.6, sub, size=14.5, colour=self.ink2, spacing=1.4)
        top, bottom = 3.0, H - 0.8
        n = max(len(steps), 1)
        gap = (bottom - top) / n
        ch = min(0.7, gap - 0.13)
        y = top
        for i, (head, note) in enumerate(steps):
            on = (highlight == i)
            self._card(s, 0.75, y, ch, ch, fill=self.brand if on else self.field[0], radius=0.5)
            self._text(s, 0.75, y, ch, ch, f"{i+1:02d}", size=11.5, font=MONO, colour="#FFFFFF",
                       bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._card(s, 0.75 + ch + 0.2, y, 11.78 - ch - 0.2, ch, fill=self.brand if on else "#FFFFFF")
            fg = "#FFFFFF" if on else self.ink
            self._text(s, 0.75 + ch + 0.55, y, 2.7, ch, head, size=13.5, font=DISPLAY,
                       bold=True, colour=fg, anchor=MSO_ANCHOR.MIDDLE)
            self._text(s, 0.75 + ch + 3.35, y, 11.78 - ch - 3.9, ch, note, size=12,
                       colour="#DCE4E8" if on else self.ink2, anchor=MSO_ANCHOR.MIDDLE)
            y += gap
        if source: self._source(s, source)
        return s

    def table(self, pill, title, sub, headers, rows, widths, source=None):
        """Brand header row, rationale column carrying the argument."""
        s = self._blank(); self._pill(s, pill)
        self._text(s, 0.75, 1.4, 11.8, 0.8, title, size=36, font=DISPLAY, bold=True, spacing=1.05)
        if sub: self._text(s, 0.75, 2.3, 11.4, 0.6, sub, size=14.5, colour=self.ink2, spacing=1.4)
        total = sum(widths); scale = 11.8 / total
        y = 3.05
        x = 0.75
        for w, hcell in zip(widths, headers):
            self._card(s, x, y, w * scale - 0.04, 0.5, fill=self.brand, radius=0.02)
            self._text(s, x + 0.18, y + 0.15, w * scale - 0.4, 0.3, hcell.upper(), size=10,
                       font=MONO, colour="#FFFFFF", bold=True)
            x += w * scale
        y += 0.58
        rh = min(0.86, (H - 1.3 - y) / max(len(rows), 1))
        for r in rows:
            x = 0.75
            self._card(s, 0.75, y, 11.76, rh - 0.06, fill="#FFFFFF")
            for w, cell in zip(widths, r):
                self._text(s, x + 0.18, y + 0.13, w * scale - 0.4, rh - 0.3, cell, size=11.5,
                           colour=self.ink2, spacing=1.3)
                x += w * scale
            y += rh
        if source: self._source(s, source)
        return s

    def statement(self, kicker, big, note=None):
        """One idea, full bleed. Used sparingly — this is the slide people remember."""
        s = self._blank(self.field[2])
        self._text(s, 1.1, 1.9, 11.2, 0.4, kicker.upper(), size=11, font=MONO, colour=self.accent, bold=True)
        self._text(s, 1.1, 2.55, 11.2, 2.6, big, size=44, font=DISPLAY, colour="#FFFFFF", bold=True, spacing=1.1)
        if note: self._text(s, 1.1, 5.35, 10, 0.9, note, size=16, colour="#CFD6DC", spacing=1.45)
        return s

    def save(self, path):
        self.p.save(path); return path
